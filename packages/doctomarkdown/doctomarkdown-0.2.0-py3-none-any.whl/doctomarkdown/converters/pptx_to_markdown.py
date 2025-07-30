from doctomarkdown.base import BaseConverter, PageResult, ConversionResult
from pptx import Presentation
import tempfile
import os
import shutil

class PptxToMarkdown(BaseConverter):
    """Converter for PPTX files to Markdown format."""
    
    def __init__(self, filepath, extract_images=False, extract_tables=False, output_path=None, llm_client=None, llm_model=None, output_type='markdown', **kwargs):
        super().__init__(filepath=filepath, extract_images=extract_images, extract_tables=extract_tables, output_path=output_path, llm_client=llm_client, llm_model=llm_model, output_type=output_type, **kwargs)

    def extract_content(self):
        temp_dir = None
        try:
            print(f"[INFO] Starting PPTX extraction for: {self.filepath}")
            # If llm_client is True, try to convert PPTX to PDF and use PDF-to-Markdown logic
            if getattr(self, 'llm_client', False):
                from doctomarkdown.utils.pptx_to_pdf import convert_pptx_to_pdf
                temp_dir = tempfile.mkdtemp()
                input_dir = os.path.dirname(self.filepath)
                output_dir = temp_dir
    
                convert_pptx_to_pdf(input_dir, output_dir)
                pptx_basename = os.path.splitext(os.path.basename(self.filepath))[0]
                pdf_path = os.path.join(output_dir, pptx_basename + '.pdf')
                if os.path.exists(pdf_path):
                    from doctomarkdown.converters.pdf_to_markdown import PdfToMarkdown
                    pdf_converter = PdfToMarkdown(pdf_path, llm_client=getattr(self, 'llm_client', None), 
                                                  llm_model=getattr(self, 'llm_model', None))
                    pages = pdf_converter.extract_content()
                    self._markdown = getattr(pdf_converter, '_markdown', None)
                    print(f"[SUCCESS] Extraction is successful via LLM for file: {self.filepath}")
                    return pages
                else:
                    print(f"[ERROR] PDF file not found after PPTX to PDF conversion: {pdf_path}")
        except Exception as e:
            print(f"[FAILURE] Error during PPTX to PDF conversion for {self.filepath}: {e}")
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
        try:
            doc = Presentation(self.filepath)
            pages = []
            markdown_lines = []
            for page_number, slide in enumerate(doc.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                page_content = "\n".join(slide_content)
                pages.append(PageResult(page_number, page_content))
                page_content = f"Page Number: {page_number}\nPage Content:\n{page_content}\n"
                markdown_lines.append(page_content)
            self._markdown = "\n".join(markdown_lines)
            print(f"[SUCCESS] Extraction is successful via standard PPTX logic for file: {self.filepath}")
            return pages
        except Exception as e:
            print(f"[FAILURE] Error during standard PPTX extraction for {self.filepath}: {e}")
            raise